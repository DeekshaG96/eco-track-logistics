from pptx import Presentation
import os

pptx_path = r"c:\Users\DEEKSHA\GOOGLESOL\eco-track-logistics\docs\Eco-Track-AI-Submission-Deck.pptx"

if not os.path.exists(pptx_path):
    print(f"Error: {pptx_path} not found.")
    exit(1)

prs = Presentation(pptx_path)

# Iterate backwards as the user mentioned the final slide (Slide 10)
# We will check all slides just to be thorough and safe.
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in paragraph.runs)
            
            # Identify if this paragraph contains placeholders
            if any(p in full_text for p in [
                "Add your Firebase Hosting URL here", 
                "Firebase Web App URL",
                "Add your YouTube or Loom video link here",
                "[Insert Live Link]",
                "[Insert Video Link]",
                "MVP Link:",
                "Working Prototype Link:",
                "Demo Video Link:"
            ]):
                
                # Perform String replacement on the aggregated text
                # We do literal replacement for all known variants
                new_text = full_text.replace("Add your Firebase Hosting URL here", "https://techspire-13303696-1c68d.web.app") \
                                    .replace("Firebase Web App URL", "https://techspire-13303696-1c68d.web.app") \
                                    .replace("[Insert Live Link]", "https://techspire-13303696-1c68d.web.app") \
                                    .replace("Add your YouTube or Loom video link here", "https://github.com/DeekshaG96/eco-track-logistics/blob/main/demo/ecotrack_final_demo.webp") \
                                    .replace("[Insert Video Link]", "https://github.com/DeekshaG96/eco-track-logistics/blob/main/demo/ecotrack_final_demo.webp")
                
                # Replace MVP/Working fields specifically if they exist with no links
                if "MVP Link:" in new_text and "http" not in new_text:
                    new_text = "MVP Link: https://techspire-13303696-1c68d.web.app"
                if "Working Prototype Link:" in new_text and "http" not in new_text:
                    new_text = "Working Prototype Link: https://techspire-13303696-1c68d.web.app"
                if "Demo Video Link:" in new_text and "http" not in new_text:
                    new_text = "Demo Video Link: https://github.com/DeekshaG96/eco-track-logistics/blob/main/demo/ecotrack_final_demo.webp"

                # Clear all runs except the first, and inject the mapped text
                if len(paragraph.runs) > 0:
                    for _ in range(1, len(paragraph.runs)):
                        p = paragraph._p
                        p.remove(paragraph.runs[1]._r)
                    paragraph.runs[0].text = new_text

prs.save(pptx_path)
print("WINNER PROTOCOL SUCCESS: PPTX Programmatically Updated via python-pptx.")
