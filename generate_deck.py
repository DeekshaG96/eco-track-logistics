from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(title, body, is_title=False):
    layout = prs.slide_layouts[0] if is_title else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    if hasattr(title_shape.text_frame.paragraphs[0], 'font'):
        title_shape.text_frame.paragraphs[0].font.bold = True
    
    if len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        body_shape.top = Inches(2.0)
        body_shape.text = body
        # Clean formatting to prevent overlapping
        for p in body_shape.text_frame.paragraphs:
            p.space_after = Pt(12)
            p.font.size = Pt(20)
    return slide

add_slide("Eco-Track AI", "Dynamic Supply Chain Optimization\nSolution Challenge 2026", True)
add_slide("1. Problem Statement", "Global logistics lacks real-time, actionable carbon visibility.\nSupply chain teams need dynamic tools to balance delivery speed with CO2 emissions.")
add_slide("2. Solution Overview", "Eco-Track AI instantly maps carbon footprints across multimodal transport options.\nEmpowering dispatchers to select the \"Greenest Choice\" operationally.")
add_slide("3. Technical Architecture", "- Frontend: React.js, Tailwind CSS, Lucide Icons\n- AI Core: Gemini 2.5 Flash\n- Hosting: Google Cloud Firebase")
add_slide("4. Generative AI Integration", "Gemini 2.5 Flash analyzes payload weight, distance, and transit nodes.\nIt returns structured JSON grading each mode's exact CO2 output and transit time.")
add_slide("5. Zero-Crash Fallback Engine", "If the Gemini AI is unreachable, the system instantly hot-swaps to an Edge Fallback.\nContinuous uptime for enterprise logistics.")
add_slide("6. Google Cloud Logistics", "Deployed via Firebase Hosting for high-availability.\nEdge-cached assets ensuring millisecond response times globally.")
add_slide("7. Sustainability Impact", "Directly targets SDG 13: Climate Action.\nEnables a projected 15-30% reduction in fleet emissions by visualizing drone and EV alternatives.")
add_slide("8. Enterprise Roadmap", "Pillar 1: Immutable Carbon Audit Trails (Smart Contracts)\nAnchoring Gemini 2.5 Flash predictions to a Layer-2 blockchain for tamper-proof ESG reporting and regulatory compliance.\n\nPillar 2: Fleet Companion App (Flutter)\nA cross-platform mobile app feeding live IoT weigh-station telemetry and GPS data directly into the React routing dashboard.")
add_slide("9. Team", "Lead Developer & Architect: Deeksha G\nTrack: Smart Supply Chains")

add_slide("10. Prototype Links", "GitHub Public Repository: https://github.com/DeekshaG96/eco-track-logistics\n\nMVP / Working Prototype Link: https://techspire-13303696-1c68d.web.app\n\nDemo Video Link: https://github.com/DeekshaG96/eco-track-logistics/blob/main/demo/ecotrack_final_demo.webp")

import os
os.makedirs(r"c:\Users\DEEKSHA\GOOGLESOL\eco-track-logistics\docs", exist_ok=True)
prs.save(r"c:\Users\DEEKSHA\GOOGLESOL\eco-track-logistics\docs\Eco-Track-AI-Submission-Deck.pptx")
print("PPTX GENERATED SUCCESSFULLY")
